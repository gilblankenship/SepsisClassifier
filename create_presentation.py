"""Generate SepsisDx PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_LIGHT = RGBColor(0x1A, 0x2D, 0x4A)
MEDICAL_BLUE = RGBColor(0x25, 0x63, 0xEB)
MEDICAL_DARK = RGBColor(0x1D, 0x4E, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_GRAY = RGBColor(0x4B, 0x55, 0x63)
MEDIUM_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
EMERALD = RGBColor(0x05, 0x96, 0x69)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_navy_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_light_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY


def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_bullet_slide(slide, items, left, top, width, font_size=16, color=DARK_GRAY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_placeholder_box(slide, left, top, width, height, label="Screenshot Placeholder"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.color.rgb = MEDIUM_GRAY
    run.font.italic = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_section_header(slide, number, title):
    add_navy_bg(slide)
    # Section number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(2.2), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MEDICAL_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(1.2),
                title, font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.5), Inches(2.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = MEDICAL_BLUE
    line.line.fill.background()


def add_page_header(slide, title, subtitle=""):
    add_white_bg(slide)
    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MEDICAL_BLUE
    bar.line.fill.background()
    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(9), Inches(0.6),
                title, font_size=28, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.85), Inches(9), Inches(0.4),
                    subtitle, font_size=14, color=MEDIUM_GRAY)
    # Brand badge
    add_textbox(slide, Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.4),
                "SepsisDx", font_size=16, color=MEDICAL_BLUE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(10.5), Inches(0.7), Inches(2.5), Inches(0.3),
                "by Digital Biosciences", font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


# =============================================================================
# SLIDE 1: Title Slide
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_navy_bg(slide)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(1.5), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = MEDICAL_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
            "SepsisDx", font_size=56, color=WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "ML-Powered Sepsis Classification from Urine Biomarkers",
            font_size=24, color=RGBColor(0x93, 0xC5, 0xFD))
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
            "Web Application Overview", font_size=18, color=MEDIUM_GRAY)

# Bottom branding
add_textbox(slide, Inches(1), Inches(6.2), Inches(5), Inches(0.4),
            "Digital Biosciences  |  An MDC Studio Portfolio Company",
            font_size=12, color=MEDIUM_GRAY)
add_textbox(slide, Inches(8), Inches(6.2), Inches(4.3), Inches(0.4),
            "sepsis-dx.com", font_size=12, color=RGBColor(0x60, 0xA5, 0xFA),
            alignment=PP_ALIGN.RIGHT)

# =============================================================================
# SLIDE 2: Executive Summary
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Executive Summary", "What is SepsisDx?")

add_bullet_slide(slide, [
    "SepsisDx is a web-based clinical decision support tool for early sepsis detection",
    "Uses machine learning to classify sepsis from urine ELISA biomarker measurements",
    "Two key biomarkers: CRP/Creatinine ratio and IP-10/Creatinine ratio (pg/mg)",
    "Validated against gold standard: Blood culture + Sepsis-3 criteria (SOFA ≥ 2)",
    "Compares ML predictions with traditional SepsisDx rule-based thresholds",
    "Five classifier options: Baseline, Logistic Regression, Random Forest, XGBoost, Neural Network",
    "Full SHAP explainability for transparent, interpretable predictions",
], Inches(0.8), Inches(1.5), Inches(7), font_size=16)

# Key stats boxes
for i, (label, value, color) in enumerate([
    ("Classifiers", "5", MEDICAL_BLUE),
    ("Biomarkers", "2", EMERALD),
    ("Explainability", "SHAP", AMBER),
    ("API", "REST", NAVY_LIGHT),
]):
    x = Inches(8.8) + Inches(0) * 0
    y = Inches(1.6) + Inches(1.3) * i
    box = add_rounded_rect(slide, Inches(9), y, Inches(3.5), Inches(1), color)
    add_textbox(slide, Inches(9.2), y + Inches(0.1), Inches(3), Inches(0.5),
                value, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.2), y + Inches(0.55), Inches(3), Inches(0.3),
                label, font_size=12, color=RGBColor(0xD1, 0xD5, 0xDB), alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 3: Architecture Overview
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Application Architecture", "Technology stack and deployment")

# Architecture boxes - left column
components = [
    ("Frontend", "Tailwind CSS, Chart.js, Bootstrap Icons\nResponsive design, dark navy theme", MEDICAL_BLUE, Inches(0.8)),
    ("Backend", "Flask + Blueprints, Gunicorn (gthread)\nSession-based auth, CORS enabled", NAVY_LIGHT, Inches(2.7)),
    ("ML Engine", "scikit-learn, SHAP, XGBoost\nBackground training with threading", EMERALD, Inches(4.6)),
]
for label, desc, color, y in components:
    add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(1.4), color)
    add_textbox(slide, Inches(1.1), y + Inches(0.15), Inches(5), Inches(0.4),
                label, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.6), Inches(5), Inches(0.7),
                desc, font_size=13, color=RGBColor(0xD1, 0xD5, 0xDB))

# Right column - deployment
add_rounded_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5), RGBColor(0xF9, 0xFA, 0xFB),
                 font_color=NAVY)
add_textbox(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.4),
            "Deployment Stack", font_size=20, color=NAVY, bold=True)

deploy_items = [
    "NVIDIA Jetson AGX Thor (aarch64)",
    "Gunicorn on port 8082",
    "Nginx reverse proxy on port 8443",
    "Cloudflare Tunnel → sepsis-dx.com",
    "Systemd service management",
    "Session auth (werkzeug hashing)",
    "Log rotation with logrotate",
]
add_bullet_slide(slide, deploy_items, Inches(7.3), Inches(2.2), Inches(5), font_size=14, color=DARK_GRAY)

# =============================================================================
# SLIDE 4: Navigation & Pages Overview
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Application Pages", "8 main sections accessible from the navigation bar")

pages = [
    ("Dashboard", "/", "bi-speedometer2", "System overview, quick predict,\nmodel stats, reference samples", MEDICAL_BLUE),
    ("Predict", "/predict/", "bi-clipboard", "Single patient classification\nwith SHAP explanations", EMERALD),
    ("Batch", "/batch/", "bi-file-spreadsheet", "CSV upload for bulk\nclassification & download", RGBColor(0x7C, 0x3A, 0xED)),
    ("Train", "/train/", "bi-gear", "Configure & run model\ntraining with CV metrics", AMBER),
    ("Models", "/models/", "bi-box", "View, activate, delete\nexported ML models", RGBColor(0xDB, 0x27, 0x77)),
    ("Explain", "/explain/", "bi-lightbulb", "SHAP plots, ROC curves,\ndecision boundaries", RGBColor(0x05, 0x96, 0x69)),
    ("Datasets", "/datasets/", "bi-database", "Browse data, generate\nsynthetic datasets", RGBColor(0x06, 0x69, 0xA2)),
    ("Help", "/help/", "bi-question", "User manual, API docs,\nclinical glossary", DARK_GRAY),
]

for i, (name, url, icon, desc, color) in enumerate(pages):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + Inches(3.15) * col
    y = Inches(1.5) + Inches(2.8) * row

    card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.4), WHITE)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)

    # Color accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.9), Pt(5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.4),
                name, font_size=20, color=NAVY, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.75), Inches(2.5), Inches(0.3),
                url, font_size=11, color=MEDICAL_BLUE)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.5), Inches(1),
                desc, font_size=13, color=DARK_GRAY)

# =============================================================================
# SLIDE 5: Login Page
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "Authentication")

# =============================================================================
# SLIDE 6: Login Detail
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Login Screen", "/auth/login — Session-based authentication")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                    "[ Screenshot: Login Page ]\n\nCentered card on navy gradient\nSepsisDx logo + Digital Biosciences branding\nUsername & password fields\nSign In button")

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
            "Features", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Standalone page (no navbar) with navy gradient background",
    "SepsisDx branding with Digital Biosciences subtitle",
    "Username & password form with icon-prefixed inputs",
    "Flash message support for error/success feedback",
    "Session-based auth using werkzeug password hashing",
    "8-hour persistent sessions",
    "All routes protected except /api/health",
    "Footer: © 2026 Digital Biosciences, an MDC Studio company",
], Inches(7.5), Inches(2.1), Inches(5), font_size=14)

# =============================================================================
# SLIDE 7: Dashboard Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "Dashboard")

# =============================================================================
# SLIDE 8: Dashboard Detail
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Dashboard", "/ — System overview and quick access to all features")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                    "[ Screenshot: Dashboard ]\n\nQuick Predict form, model stats cards,\nnavigation grid, reference samples table")

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
            "Functions", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Quick Predict form — enter CRP/Cr and IP-10/Cr for instant classification",
    "System statistics cards: total models, best AUROC, active model",
    "6-card navigation grid linking to all major features",
    "Reference samples table with 6 pre-configured test cases",
    "Highlights false negative cases (Sample 6) for awareness",
    "Shows both ML prediction and SepsisDx rule-based comparison",
    "Active model badge in navbar with real-time indicator",
], Inches(7.5), Inches(2.1), Inches(5), font_size=14)

# =============================================================================
# SLIDE 9: Predict Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "Single Prediction")

# =============================================================================
# SLIDE 10: Predict Form
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Prediction — Input Form", "/predict/ — Enter biomarker values for classification")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                    "[ Screenshot: Predict Form ]\n\nBiomarker input fields,\nmodel selector dropdown,\nreference table, risk legend")

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
            "Input Form Features", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "CRP/Creatinine ratio input (pg/mg) — threshold hint: >300 = Septic",
    "IP-10/Creatinine ratio input (pg/mg) — threshold hint: >100 = Septic",
    "Model selection dropdown showing all models with AUROC scores",
    "Biomarker reference table with clinical thresholds",
    "Risk level color legend: LOW, MODERATE, HIGH, VERY HIGH",
    "Validates inputs before submission",
], Inches(7.5), Inches(2.1), Inches(5), font_size=14)

# =============================================================================
# SLIDE 11: Predict Result
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Prediction — Results", "/predict/run — Classification result with SHAP explanation")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                    "[ Screenshot: Prediction Result ]\n\nColor-coded result card (red/green),\nprobability bar, SHAP table,\nclinical notes")

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
            "Result Features", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Large color-coded result: Red (Septic) or Green (Not Septic)",
    "Probability percentage with visual progress bar",
    "Risk level badge: LOW / MODERATE / HIGH / VERY HIGH",
    "Optimal threshold indicator used for classification",
    "SepsisDx rule-based comparison (agreement check)",
    "SHAP contribution table per feature with direction arrows",
    "Clinical context note with actionable guidance",
    "Links to new prediction and full SHAP dashboard",
], Inches(7.5), Inches(2.1), Inches(5), font_size=14)

# =============================================================================
# SLIDE 12: Batch Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "Batch Processing")

# =============================================================================
# SLIDE 13: Batch Upload & Results
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Batch Classification", "/batch/ — Upload CSV for bulk predictions")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.3),
                    "[ Screenshot: Batch Upload ]\nCSV upload form with format guide")
add_placeholder_box(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.6),
                    "[ Screenshot: Batch Results ]\nSummary stats, results table, download button")

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Upload", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "CSV file upload with format specification",
    "Required columns: crp_cr_pg_mg, ip10_cr_pg_mg",
    "Optional: sample_id, clinical_status, creatinine",
    "Example CSV format displayed for reference",
], Inches(7), Inches(2.1), Inches(5.5), font_size=14)

add_textbox(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.4),
            "Results", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Summary cards: total, septic, not septic, agreement %",
    "Results table showing first 100 rows",
    "Full CSV download with all predictions",
    "Columns: probability, risk level, SepsisDx comparison",
], Inches(7), Inches(4.7), Inches(5.5), font_size=14)

# =============================================================================
# SLIDE 14: Train Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "Model Training")

# =============================================================================
# SLIDE 15: Train Configure
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Training — Configuration", "/train/ — Select classifiers, dataset, and parameters")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                    "[ Screenshot: Training Configuration ]\n\nClassifier checkboxes with parameters,\ndataset selection, feature mode")

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
            "Configuration Options", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "5 classifiers available:",
    "  • SepsisDx Baseline (rule-based threshold)",
    "  • Logistic Regression (C, class_weight)",
    "  • Random Forest (n_estimators, max_depth)",
    "  • XGBoost (learning_rate, max_depth)",
    "  • Neural Network (hidden layers, max_iter)",
    "Per-classifier parameter customization",
    "Dataset selection: existing CSVs or file upload",
    "Feature mode: Basic (2 features) or Extended (7)",
    "Background training with real-time progress",
], Inches(7.5), Inches(2.1), Inches(5), font_size=13)

# =============================================================================
# SLIDE 16: Train Progress & Results
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Training — Progress & Results", "Real-time progress tracking and cross-validation metrics")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.3),
                    "[ Screenshot: Training Progress ]\nSpinner, progress bar, elapsed time")
add_placeholder_box(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.6),
                    "[ Screenshot: Training Results ]\nCV metrics table, AUROC chart, exported models")

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Progress Tracking", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Animated spinner with status messages",
    "Progress bar updates every 2-3 seconds",
    "Elapsed time counter",
    "Auto-redirects to results on completion",
], Inches(7), Inches(2.1), Inches(5.5), font_size=14)

add_textbox(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.4),
            "Results Display", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "5-fold cross-validation metrics table",
    "AUROC, Sensitivity, Specificity, F1 (mean ± std)",
    "AUROC comparison bar chart",
    "List of exported model files",
    "Links to Models page and SHAP analysis",
], Inches(7), Inches(4.7), Inches(5.5), font_size=14)

# =============================================================================
# SLIDE 17: Models Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "Model Management")

# =============================================================================
# SLIDE 18: Models List & Detail
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Model Management", "/models/ — View, activate, and manage trained models")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.3),
                    "[ Screenshot: Models List ]\nGrid of model cards with active badge")
add_placeholder_box(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.6),
                    "[ Screenshot: Model Detail ]\nMetadata, parameters, training metrics")

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Models List", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Responsive card grid (1-4 columns)",
    "Active model highlighted with colored ring",
    "Shows: name, classifier type, AUROC, threshold, export date",
    "Actions: View Details, Activate, Delete (with confirmation)",
], Inches(7), Inches(2.1), Inches(5.5), font_size=14)

add_textbox(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.4),
            "Model Detail", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Full metadata: name, class, threshold, export time",
    "Gold standard and biomarker source info",
    "Model hyperparameters table",
    "Training metrics (AUROC, sensitivity, specificity, etc.)",
    "File size information",
], Inches(7), Inches(4.7), Inches(5.5), font_size=14)

# =============================================================================
# SLIDE 19: Explain Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 7, "Explainability (SHAP)")

# =============================================================================
# SLIDE 20: Explain Dashboard
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "SHAP Explainability Dashboard", "/explain/ — Model interpretation and visualization")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                    "[ Screenshot: Explain Dashboard ]\n\nTabbed interface with 5 visualization types:\nSHAP Summary, Feature Importance,\nDecision Boundary, ROC Curves, Calibration")

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
            "Visualization Tabs", font_size=20, color=NAVY, bold=True)

viz_items = [
    "SHAP Summary Plot — Feature impact distribution across all samples",
    "Feature Importance — Mean |SHAP| values ranked by contribution",
    "Decision Boundary — 2D visualization of classification regions",
    "ROC Curves — Model performance comparison with literature benchmarks",
    "Calibration Curves — Predicted vs. actual probability reliability",
]
add_bullet_slide(slide, viz_items, Inches(7.5), Inches(2.1), Inches(5), font_size=14)

add_textbox(slide, Inches(7.5), Inches(4.8), Inches(5), Inches(0.4),
            "Key Features", font_size=18, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Model selector dropdown for comparison",
    "Plots generated on-demand (server-side)",
    "Lazy-loaded images for performance",
    "Dynamic updates when switching models",
], Inches(7.5), Inches(5.3), Inches(5), font_size=13)

# =============================================================================
# SLIDE 21: Datasets Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 8, "Dataset Management")

# =============================================================================
# SLIDE 22: Datasets
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Dataset Management", "/datasets/ — Browse data and generate synthetic datasets")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.3),
                    "[ Screenshot: Datasets List ]\nDataset table + synthetic generator form")
add_placeholder_box(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.6),
                    "[ Screenshot: Dataset Viewer ]\nStats, class distribution, paginated data table")

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Dataset Browser", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Table: dataset name, rows, columns, file size, download",
    "Click to open dataset viewer with full data preview",
], Inches(7), Inches(2.1), Inches(5.5), font_size=14)

add_textbox(slide, Inches(7), Inches(3.2), Inches(5.5), Inches(0.4),
            "Synthetic Data Generator", font_size=18, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Configure: septic, SIRS, healthy sample counts",
    "Population profiles: General, ICU, Pediatric",
    "Noise levels: Low, Medium, High",
    "Optional extended biomarkers (7 features)",
], Inches(7), Inches(3.7), Inches(5.5), font_size=14)

add_textbox(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(0.4),
            "Dataset Viewer", font_size=18, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Column statistics (mean, std, min, max)",
    "Class distribution badges (color-coded)",
    "Paginated data table (50 rows/page)",
    "CSV download for any dataset",
], Inches(7), Inches(5.5), Inches(5.5), font_size=14)

# =============================================================================
# SLIDE 23: Help Section Header
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 9, "Help & API")

# =============================================================================
# SLIDE 24: Help & API
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_page_header(slide, "Help Page & REST API", "/help/ and /api/ — Documentation and programmatic access")

add_placeholder_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2),
                    "[ Screenshot: Help Page ]\n\nTable of contents, expandable sections,\nbiomarker reference, classifier guide,\nAPI docs, clinical glossary")

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Help Page (14 Sections)", font_size=20, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "Introduction & Quick Start guide",
    "Biomarker reference (CRP/Cr, IP-10/Cr thresholds)",
    "Classifier descriptions and use cases",
    "Using predictions, risk levels, clinical notes",
    "Training pipeline and SHAP interpretation",
    "Input data format with example CSV",
    "Threshold optimization strategies",
    "Clinical & ML glossary (15 terms)",
], Inches(7), Inches(2.1), Inches(5.5), font_size=14)

add_textbox(slide, Inches(7), Inches(4.8), Inches(5.5), Inches(0.4),
            "REST API Endpoints", font_size=18, color=NAVY, bold=True)
add_bullet_slide(slide, [
    "GET  /api/health — Service health check (no auth)",
    "POST /api/predict — JSON prediction (crp_cr, ip10_cr)",
    "GET  /api/models — List all exported models",
], Inches(7), Inches(5.3), Inches(5.5), font_size=14)

# =============================================================================
# SLIDE 25: Summary / Closing
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_navy_bg(slide)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
            "SepsisDx — Complete ML Pipeline for Sepsis Detection",
            font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Feature summary grid
features = [
    ("Predict", "Single & batch patient\nclassification with\nprobability scoring"),
    ("Train", "5 classifiers with\ncross-validation and\nthreshold optimization"),
    ("Explain", "SHAP plots, ROC curves,\ndecision boundaries\nfor transparency"),
    ("Manage", "Model versioning,\nactivation, dataset\nbrowsing & generation"),
]
for i, (title, desc) in enumerate(features):
    x = Inches(0.8) + Inches(3.15) * i
    add_rounded_rect(slide, x, Inches(2.8), Inches(2.8), Inches(2.3), NAVY_LIGHT)
    add_textbox(slide, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(0.4),
                title, font_size=22, color=MEDICAL_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(1.5),
                desc, font_size=14, color=RGBColor(0xD1, 0xD5, 0xDB), alignment=PP_ALIGN.CENTER)

# Bottom branding
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.6), Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = MEDICAL_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.9), Inches(11.3), Inches(0.5),
            "Digital Biosciences  •  An MDC Studio Portfolio Company",
            font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.4), Inches(11.3), Inches(0.4),
            "sepsis-dx.com  |  dbmsciences.com  |  mdcstudio.com",
            font_size=12, color=RGBColor(0x60, 0xA5, 0xFA), alignment=PP_ALIGN.CENTER)

# Save
output_path = "SepsisDx_Application_Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
